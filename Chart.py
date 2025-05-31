from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from pptx.dml.color import RGBColor

# Create a new colorful presentation
prs = Presentation()

# Define a colorful background title slide
title_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(title_slide_layout)
title_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
tf = title_box.text_frame
title_run = tf.add_paragraph().add_run()
title_run.text = "100-Day Roadmap: Data/Analytics Engineer Growth Plan"
title_run.font.size = Pt(36)
title_run.font.bold = True
title_run.font.color.rgb = RGBColor(255, 255, 255)
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = RGBColor(0, 102, 204)

subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
tf_sub = subtitle_box.text_frame
sub_run = tf_sub.add_paragraph().add_run()
sub_run.text = "Includes dbt Certification, CI/CD, Airflow, and LeetCode Integration"
sub_run.font.size = Pt(20)
sub_run.font.color.rgb = RGBColor(255, 255, 255)

# Define a color palette for phases
phase_colors = [
    RGBColor(255, 204, 0),    # Yellow
    RGBColor(102, 204, 0),    # Green
    RGBColor(0, 153, 204),    # Blue
    RGBColor(204, 102, 255)   # Purple
]

# Add slides with color-coded title bars and checkboxes
for idx, phase in enumerate(phases):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    color = phase_colors[idx % len(phase_colors)]

    # Colored title bar
    title_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1)
    )
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = color
    title_shape.text = phase["title"]
    title_shape.text_frame.paragraphs[0].font.size = Pt(24)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # Add checklist items
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    for task in phase["tasks"]:
        p = tf.add_paragraph()
        p.text = f"☐ {task}"
        p.level = 0
        p.font.size = Pt(18)

# Save colorful version
pptx_path_color = "/mnt/data/100_Day_Data_Engineer_Roadmap_Colorful.pptx"
prs.save(pptx_path_color)
pptx_path_color
